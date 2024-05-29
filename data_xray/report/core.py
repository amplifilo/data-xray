from ..scan import PlotImage
from ..grid import *
#from ..grid.dsops import grid_viz
from ..nanonisio import Grid, Scan
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from matplotlib import pyplot as plt
import pandas as pd

#some utilies to work with powerpoint
meanstd = lambda arr: np.std(np.ravel(arr)/np.mean(np.ravel(arr)))


class SummaryPPT(object):
    #group-sequences selects sequential images with group_sequences images of the same area (like sequential slides)
    def __init__(self, pname="image_summary", new=False, fdict=None, group_sequences = False, maximages=50, 
                    inject_method="insert_data", chanselect = {'scan':'Z', 'grid':'cf'}, **kwargs):
    
        if fdict is None:
            print('please specify data to summarize')
            return
        elif group_sequences:
            fdicts = []
            scan_offsets  = [tuple(f.header['scan_offset']) for f in fdict]
            pddf = pd.DataFrame({'offset':scan_offsets, 'fname':fdict});
            pddf.groupby('offset').ngroups
            for k, v in pddf.groupby('offset').groups.items():
                if len(v) > group_sequences: #group_sequences can be as small as 1 (for 2 sequences)
                    fdicts.append([fdict[vv] for vv in v])            
            print('movies found')
            print([len(f) for f in fdicts])
            
        else:
            fdicts = [fdict[i:i + maximages] for i in range(0, len(fdict), maximages)]

        self.topdir = os.path.commonpath([j.fname for j in fdict]) + '/'
        self.chanselect = chanselect
        for j,f in enumerate(fdicts): 
            self.presentation_name = pname + '_' + str(j)
            self.pptx_file_name = self.topdir + self.presentation_name + '.pptx'
            self.fdict = f
            #self.new = new
            self.init_ppt(self.presentation_name)
            
            #self.insert_images()
            insert_func = getattr(self, inject_method, None)
            insert_func()
            #try:
            self.pres.save(self.pptx_file_name)
            print('batch ' + str(j) +' stored in : ' + self.pptx_file_name )
            #except:
            #    print('something wrong with saving the presentation file ' + self.pptx_file_name)

    def init_ppt(self, presentation_name):
        pres = Presentation()
        pres.notes_master.name = self.presentation_name
        self.pres = pres

    def insert_data(self):
        for fj in self.fdict:
            if re.findall('sxm', fj.fname):
                    s_d = Scan(fj.header['fname'], header_only=False)
                #try:
                    if self.chanselect == "Automatic":
                        # attempt to recognize good data
                        plotsignals = []
                        for c in s_d.signals.keys():
                            sig = s_d.signals[c]['forward']
                            sig = sig[~np.isnan(sig)]
                            if meanstd(sig) > 2:
                                plotsignals.append(c)
                        plotsignals.append('Z')
                        plotsignals = list(set(plotsignals))  # no dobule Z

                    else:
                        plotsignals = self.chanselect['scan']
                    
                    plotsignals = [p for p in plotsignals if p in s_d.signals.keys()]

                    nrows = 2 if len(plotsignals) > 3 else 1
                    ncols = int(np.ceil(len(plotsignals) / nrows))
                    #f3, a3 = plt.subplots(nrows, ncols);
                    f3, a3 = pplt.subplots(nrows=nrows, ncols=ncols);
                    if a3 is not list:
                        a3 = [a3]
                    for c, a in zip(plotsignals, np.ravel(a3)):
                        PlotImage(s_d, chan=c, ax=a, high_pass=None);

                    [a.axis('off') for a in np.ravel(a3)]

                    xy = ['X', 'Y']
                    offset = s_d.header['scan_offset'] / 1e-9
                    xyoffsets = [xy[j] + '=' + str(np.round(offset[j], 2)) + ' nm ' for j in [0, 1]]

                    titleString = [s_d.fname]
                    titleString.append('Bias: ' + str(s_d.header['bias']) + 'V')
                    titleString.append('Control: ' + s_d.header['z-controller']['Name'][0] + ' Setpoint:' + str(s_d.header['z-controller']['Setpoint'].values)[2:-2])
                    titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
                    titleString.append('Resolution: ' + str(s_d.header['scan_pixels']))

                    self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                    print(os.path.basename(s_d.ds.fname) + ' imported')
                    pplt.close();
                    del s_d
                    #f3.clf();  # close figure so that it doesn't clog up in the end
                #except:
                #    print(os.path.basename(fj.fname) + ' failed')

            elif re.findall('3ds', fj.fname):
                #fix the same way as images. load each file name, then delete the object
                try:
                    g_d = Grid(fj.fname)

                    fig, ax = pplt.subplots(nrows = 1, ncols = 2);

                    grid_viz.chan_histogram(g_d.ds, xy=['bias', self.chanselect['grid']], xymod=[lambda x: x, lambda x: x / np.mean(np.ravel(x))],
                                    ax=ax[0], label=['bias', self.chanselect['grid'], ''])

                    # plop in a clustered map
                    km = grid_viz.chan_pca_kmeans(g_d.ds, xvec='bias', chan=self.chanselect['grid'], mod=lambda x: x / np.mean(np.ravel(x)),
                                         comps=6, nclust=4, fig=None);
                    ax[1].imshow(km);
                    
                    titleString = [fj.fname]
                    titleString.append('Pixels:' + str(fj.ds.attrs['dim_px']))
                    titleString.append('Size:' + str(fj.ds.attrs['size_xy'][0])+'/'+str(fj.ds.attrs['size_xy'][1]))
                    titleString.append('Date:' + str(fj.ds.attrs['start_time']))
                
                    self.fig_to_ppt([fig], leftop=[1, 2], txt=titleString)
                    print(os.path.basename(fj.fname) + ' imported')

                    pplt.close();  # close figure so that it doesn't clog up in the end
                    del fig
                    #fig.clf();
                    del g_d

                except:
                    print(os.path.basename(fj.fname) + ' import into ppt failed')
                    continue

    def insert_maps(self, mapchan=['cf']):
        
        #plop in a histogram of the map channel across the whole thing
        
        for fj in self.fdict:
            fig,ax=plt.subplots(2,1, figsize=(8,5))
            ChanHistogramDS(fj.ds, xy=['bias',mapchan[0]], xymod=[lambda x:x,lambda x:x/np.mean(np.ravel(x))], ax=ax[0], label=['bias',mapchan[0][:-1],''])  
            
            #plop in a clustered map
            km = ChanPcaKmeansDS(fj.ds, xvec='bias', chan=mapchan[0], mod = lambda x: x/np.mean(np.ravel(x)), comps=6, nclust=4, fig=None)
            ax[1].imshow(km)
            plt.colorbar()
        
            #fig.savefig("pca_3d.png",bbox_inches='tight')

            ###need to add name attribute to grids
            titleString = [fj.fname]
        

            # titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
            # titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
            # titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
            # titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
            try:
                self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                print(os.path.basename(fj.ds.fname) + ' imported')
            
                f3.clf(); #close figure so that it doesn't clog up in the end
            
            except:
                print(os.path.basename(fj.ds.fname) + ' failed')
            
            #self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)

        return

    def insert_Z_images(self):
        """
        Dumpt a batch of images into a powerpoint

        :param chanselect:
        :param fdict:
        :param topdir:
        :return:
        """
        #for folder in (self.fdict.keys()):
        #        self.text_to_slide(folder)

        # newpres = Presentation()
        # newpres.notes_master.name = 'sum1.pptx'
        # newpres.save(newpres.notes_master.name)

        for fj in self.fdict:
            # TextToSlide(fj.fname,pres=pres)
            try:
                # scf = {1:(12,6), 2:(12,9), 3:(12,10)}
                if self.chanselect == "Automatic":
                #attempt to recognize good data  
                    plotsignals = []
                    for c in fj.signals.keys():
                        sig = fj.signals[c]['forward']
                        sig = sig[~np.isnan(sig)]
                        if meanstd(sig) > 2:
                            plotsignals.append(c)    
                    plotsignals.append('Z')
                    plotsignals = list(set(plotsignals)) #no dobule Z
                    

                else:
                    plotsignals = self.chanselect
                    
                nrows = 2 if len(plotsignals) > 2 else 1
                ncols = int(np.ceil(len(plotsignals)/nrows))
                f3, a3 = plt.subplots(nrows,ncols);
                if a3 is not list:
                    a3 = [a3]
                for c,a in zip(plotsignals,np.ravel(a3)):
                    PlotImage(fj, chan=c, ax=a, high_pass=None);
                
                [a.axis('off') for a in np.ravel(a3)]

               

                xy = ['X','Y']
                offset = fj.header['scan_offset']/1e-9
                xyoffsets = [xy[j] + '=' + str(np.round(offset[j],2)) + ' nm ' for j in [0,1]]

                titleString = [fj.fname]
                titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
                setpoint = str(fj.ds.attrs['z-controller']['Setpoint'].values)[2:-2]
                titleString.append('Control: ' + fj.header['z-controller']['Name'][0] + ' Setpoint: ' + setpoint)
                titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
                titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
                
                self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)
                print(os.path.basename(fj.ds.fname) + ' imported')
                f3.clf(); #close figure so that it doesn't clog up in the end
            except:
                print(os.path.basename(fj.ds.fname) + ' failed')

    def insert_photos(self):
        #assumes that fdict is just filenames

        for fj in self.fdict:
            
            slide_layout = self.pres.slide_layouts[5]  # 6 is the index for a blank slide
            slide = self.pres.slides.add_slide(slide_layout)
        
            left = top = Inches(1)
            pic = slide.shapes.add_picture(fj.fname, left, top, width=self.pres.slide_width - Inches(2))
            
            print(os.path.basename(fj.fname) + ' imported')

            title = slide.shapes.title
            title.text = fj.fname
            #self.text_to_slide(fj.fname, slide=slide)

            # titleString.append('Bias: ' + str(fj.header['bias']) + 'V')
            # titleString.append('Control: ' + fj.header['z-controller']['Name'][0])
            # titleString.append('Offsets: ' + xyoffsets[0] + xyoffsets[1])
            # titleString.append('Resolution: ' + str(fj.header['scan_pixels']))
            
            #self.fig_to_ppt([f3], leftop=[1, 2], txt=titleString)

        return

    def fig_to_ppt(self, figs, leftop=[0,1.5], txt=None):
        """
        Plop figures into powerpoint
        :param figs:
        :param pres:
        :param leftop:
        :param txt:
        :return:
        """
        #savepptx needs to be a full path. If None is provided the default presentation
        #will be created with a name sum1.pptx in the current folder
        from pptx.util import Inches

        blank_slide_layout = self.pres.slide_layouts[5]
        left = Inches(leftop[0])
        top = Inches(leftop[1])

        tmp_path = 't1.png'
        for figp in figs:
            plt.savefig(tmp_path, transparent=1, format='png', dpi=300, bbox_inches = 'tight')
            slide = self.pres.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(tmp_path, left, top)

        if txt is not None:
            self.text_to_slide(txt, slide=slide)

    def text_to_slide(self, txt, slide=None): #lets make txt a list of strings
                                              
        """
        convert text to slide

        :param txt: list of strings
        :param pres:
        :param slide:
        :return:
        """
        from pptx.util import Pt
        #title = slide.shapes.title
        #subtitle = slide.placeholders[1]

       # title.text = "Hello, World!"
        #subtitle.text = "python-pptx was here!"

       # prs.save('test.pptx')

        from pptx.util import Inches

        if self.pres == None:
            print('please init presentation')
        else:
            if slide is None:
                bullet_slide_layout = self.pres.slide_layouts[5]
                slide = self.pres.slides.add_slide(bullet_slide_layout)

            shapes = slide.shapes

            countshapes = 0

            #just catch the first shape object with a frame in it
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elif countshapes > 0:
                    tframe = shape.text_frame
                    tframe.clear()
                    #print('caught one')
                else:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    countshapes = 1

            text_frame.clear()
            p = text_frame.paragraphs[0]
            
            for t in txt:
                run = p.add_run()
                run.text = t

                font = run.font
                font.name = 'Calibri'
                font.size = Pt(12)    
                #    p.text = t
                p = text_frame.add_paragraph()



class QuickPPT(object):
#similar to SummaryPPT, but now decoupled from specific 
# data types. Designed to be a catch-all, python -> ppt tool.
# Todo: update SummaryPPT to utilize QuickPPT constructs.
#     
    def __init__(self, presentation_name):
        self.presentation_name = presentation_name
        self.init_ppt()

    def init_ppt(self):
        pres = Presentation()
        pres.notes_master.name = self.presentation_name
        self.pres = pres

    def save(self):
        self.pres.save(self.presentation_name)

    def png_to_ppt(self, pngfile, ttl = []):
       """
       Plop a PNG file into powerpoint slide
       :param pngfile:
       :param pres:
       :param ttl:
       :return:
       """

       #blank_slide_layout = pres.slide_layouts[6]
       title_slide_layout = self.pres.slide_layouts[9]
       left = top = Inches(1)

       slide = self.pres.slides.add_slide(title_slide_layout)
       slide.shapes.add_picture(pngfile, left, top)
       subtitle = slide.placeholders[1]
       title = slide.shapes.title
       if len(ttl):
           subtitle.text = ttl

    def pplt_to_ppt(self, figs, leftop=[[0,1.2]], txt=None):
        """
        Plop figures into powerpoint
        :param figs:
        :param pres:
        :param leftop:
        :param txt:
        :return:
        """
        #savepptx needs to be a full path. If None is provided the default presentation
        #will be created with a name sum1.pptx in the current folder
        from pptx.util import Inches

        blank_slide_layout = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(blank_slide_layout)
        
        
        tmp_path = 't1.png'
        for figp, figoffset in zip(figs, leftop):
            left = Inches(figoffset[0])
            top = Inches(figoffset[1])
            
            figp.savefig(tmp_path, transparent=1, format='png', dpi=300)
            slide.shapes.add_picture(tmp_path, left, top)

        if txt is not None:
            self.text_to_slide(txt, slide=slide)

    def text_to_slide(self, txt, slide=None): #lets make txt a list of strings
                                              
        """
        convert text to slide

        :param txt: list of strings
        :param pres:
        :param slide:
        :return:
        """
        from pptx.util import Pt
        #title = slide.shapes.title
        #subtitle = slide.placeholders[1]

       # title.text = "Hello, World!"
        #subtitle.text = "python-pptx was here!"

       # prs.save('test.pptx')

        from pptx.util import Inches

        if self.pres == None:
            print('please init presentation')
        else:
            if slide is None:
                bullet_slide_layout = self.pres.slide_layouts[6]
                slide = self.pres.slides.add_slide(bullet_slide_layout)

            shapes = slide.shapes

            countshapes = 0

            #just catch the first shape object with a frame in it
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elif countshapes > 0:
                    tframe = shape.text_frame
                    tframe.clear()
                    #print('caught one')
                else:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    countshapes = 1

            text_frame.clear()
            p = text_frame.paragraphs[0]
            
            for t in txt:
                run = p.add_run()
                run.text = t

                font = run.font
                font.name = 'Calibri'
                font.size = Pt(14)
                #    p.text = t
                p = text_frame.add_paragraph()     